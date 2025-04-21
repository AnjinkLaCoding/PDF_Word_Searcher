# %%
import os
import fitz
import tkinter as tk
from tkinter import filedialog, messagebox, simpledialog
from pptx import Presentation
import pandas as pd
from docx import Document
import win32com.client
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from urllib.parse import quote

# %%
def extract_text_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    text = "".join(page.get_text("text") for page in doc)
    return text

# %%
def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "".join(text)

# %%
def extract_text_from_xlsx(exc_path):
    text = []
    df = pd.read_excel(exc_path)  # Read all sheets
    for sheet, data in df.items():
        temp = df.dropna(how="all").astype(str).stack().tolist()  # Drop NaNs and flatten
        text.append("\n".join(temp))
    text = "".join(text)
    return text

# %%
def convert_doc_to_docx(doc_path, docx_path):
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False  # Set to True if you want to see Word doing the conversion
    doc = word.Documents.Open(doc_path)
    doc.SaveAs(docx_path, FileFormat=16)  # FileFormat 16 = Word Document (.docx)
    doc.Close()
    word.Quit()

# %%
def extract_text_from_docx(file_path):
    if file_path.endswith(".docx"):
        doc = Document(file_path)
        text = [para.text for para in doc.paragraphs if para.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                text.append("\t".join(cell.text.strip() for cell in row.cells))
    elif file_path.endswith(".doc"):
        docx_path = file_path.replace(".doc",".docx")
        convert_doc_to_docx(file_path, docx_path)
        doc = Document(docx_path)
        text = [para.text for para in doc.paragraphs if para.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                text.append("\t".join(cell.text.strip() for cell in row.cells))
    return "\n".join(text)

def extract_text_from_log(file_path):
    text=[]
    with open(file_path, "r", encoding="utf-8") as file:
        for line in file:
            text.append(line)
    return "".join(text)

# %%
def index_pdfs():
    folder_selected = filedialog.askdirectory()
    DirList=[]
    if folder_selected:
        entry_pdf.delete(0, tk.END)
        entry_pdf.insert(0, folder_selected)
    folder_path = entry_pdf.get()
    for dirname in os.scandir(folder_selected):
        if dirname.is_dir():
            DirList.append({"DirName":dirname.name,"DirPath":dirname.path})
    for j in DirList:
        j["DirPath"] = j["DirPath"].replace("/","\\")
    for i in DirList:
        for filename in os.listdir(i["DirPath"]):
            if filename.endswith(".pdf"):
                file_path = os.path.join(i["DirPath"], filename)
                text = extract_text_from_pdf(file_path)
                file_data.append({"path": file_path, "filename": filename, "content": text})
            elif filename.endswith(".pptx"):
                file_path = os.path.join(i["DirPath"], filename)
                text = extract_text_from_pptx(file_path)
                file_data.append({"path": file_path, "filename": filename, "content": text})
            elif filename.endswith(".docx"):
                file_path = os.path.join(i["DirPath"], filename)
                text = extract_text_from_docx(file_path)
                file_data.append({"path": file_path, "filename": filename, "content": text})
            elif filename.endswith(".xlsx"):
                file_path = os.path.join(i["DirPath"], filename)
                text = extract_text_from_xlsx(file_path)
                file_data.append({"path": file_path, "filename": filename, "content": text})
            elif filename.endswith(".doc"):
                file_path = os.path.join(i["DirPath"], filename)
                text = extract_text_from_docx(file_path)
                file_data.append({"path": file_path, "filename": filename, "content":text})
            elif filename.endswith(".log"):
                file_path = os.path.join(i["DirPath"], filename)
                text = extract_text_from_log(file_path)
                file_data.append({"path": file_path, "filename": filename, "content":text})
    if file_data:
        messagebox.showinfo("Success", "Files indexed successfully!")
    else:
        messagebox.showwarning("Warning", "Error Occured!")

# %%
def save_as():
    found_files = []
    File_path = []
    query = entry_words.get()
    if not query:
        messagebox.showinfo("Results", "Please enter a word")
        return
    for pdf in file_data:
        if query.lower() in pdf["content"].lower():
            found_files+=[pdf["filename"]]
            temp = pdf["path"].replace('\\', '/')
            temp1 = f"file:///{quote(temp)}"
            File_path+=[temp1]
    if not found_files:
        messagebox.showinfo("Results", "The word can't be found!")
        return
    save_path = entry_output.get()
    pdfmetrics.registerFont(TTFont('NotoFull', "C:/Users/B103040059/Downloads/NotoSansTC-VariableFont_wght.ttf")) ##Font pack for Traditional Chinese, edit it accoridng to your path
    if save_path:
        c = canvas.Canvas(save_path, pagesize=A4)
        width, height = A4
        text_y = height - inch
        c.setFont("NotoFull", 12)
        c.drawString(100, text_y + 20, f"Click the links below to open the files with the word {query}:")
        for i in range(len(found_files)):
            c.drawString(100, text_y, found_files[i])
            c.linkURL(File_path[i], (100, text_y - 2, 300, text_y + 12), relative=0)
            text_y -= 20
        c.save()
        messagebox.showinfo("Saved", "Results saved successfully!")
    else:
        messagebox.showwarning("Warning", "Error saving file!")

# %%
def select_output():
    file_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF Files", "*.pdf")])
    if file_path:
        entry_output.delete(0, tk.END)
        entry_output.insert(0, file_path)

# %%
root = tk.Tk()
root.title("PDF Word Searcher")
root.geometry("400x250")

file_data = []
tk.Label(root, text="Select Folder:").pack()
entry_pdf = tk.Entry(root, width=50) ##variable to pass the user input to the function
entry_pdf.pack()
tk.Button(root, text="Browse", command=index_pdfs).pack()

tk.Label(root, text="Enter words to search:").pack()
entry_words = tk.Entry(root, width=50)
entry_words.pack()

tk.Label(root, text="Save As:").pack()
entry_output = tk.Entry(root, width=50)
entry_output.pack()
tk.Button(root, text="Browse", command=select_output).pack()

tk.Button(root, text="Generate List file", command=save_as).pack()
root.mainloop()