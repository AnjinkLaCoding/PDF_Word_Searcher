#!/usr/bin/env python
# coding: utf-8

# In[1]:


import os
import fitz
import tkinter as tk
from tkinter import filedialog, messagebox, simpledialog, ttk
from pptx import Presentation
import pandas as pd
from docx import Document
import win32com.client
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from urllib.parse import quote


# In[2]:


def extract_text_from_pdf(pdf_path):
    try:
        doc = fitz.open(pdf_path)
        text = "".join(page.get_text("text") for page in doc)
    except:
        ##messagebox.showwarning("Warning", f"{pdf_path} might be corrupted or unreadable")
        text=""
        error_files.append(pdf_path)
    return text


# In[3]:


def extract_text_from_pptx(file_path):
    text = []
    try:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
    except:
        ##messagebox.showwarning("Warning", f"{file_path} might be corrupted or unreadable")
        text=""
        error_files.append(file_path)
    return "".join(text)


# In[4]:


def extract_text_from_xlsx(exc_path):
    Flag=0;
    try:
        text = entry_words.get().lower()
        df = pd.read_excel(exc_path, engine='openpyxl')  # Read all sheets
        Flag = df.apply(lambda row: row.astype(str).str.contains(text, case=False).any(), axis=1).any()
    except:
        ##messagebox.showwarning("Warning", f"{exc_path} might be corrupted or unreadable")
        text=""
        error_files.append(exc_path)
    return Flag


# In[5]:


def convert_doc_to_docx(doc_path):
    word = win32com.client.Dispatch("Word.Application")
    word.Visible = False
    word.DisplayAlerts = 0  # Disable alerts/popups
    try:
        doc = word.Documents.Open(doc_path)
        new_doc = word.Documents.Add()

        docx_path = doc_path.replace(".doc", ".docx")

        doc.Content.Copy()
        new_doc.Content.Paste()

        new_doc.SaveAs(docx_path, FileFormat=16) # Save as .docx format (FileFormat=16 is for .docx)
        Flag_doc = 1
    except Exception as e:
        print(f"Error converting {doc_path}: {e}")
        Flag_doc = 0
    finally:
        if 'doc' in locals():
            doc.Close(False)
        if 'new_doc' in locals():
            new_doc.Close(False)
        word.Quit()


# In[6]:


def extract_text_from_docx(file_path):
    if file_path.endswith(".docx"):
        try:
            doc = Document(file_path)
            text = [para.text for para in doc.paragraphs if para.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    text.append("\t".join(cell.text.strip() for cell in row.cells))
        except:
            text=""
            error_files.append(file_path)
    elif file_path.endswith(".doc"):
        docx_path = file_path.replace(".doc",".docx")
        convert_doc_to_docx(file_path)
        if Flag_doc == 0:
            text=""
            error_files.append(file_path)
        else:
            doc = Document(docx_path)
            text = [para.text for para in doc.paragraphs if para.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    text.append("\t".join(cell.text.strip() for cell in row.cells))
    return "\n".join(text)


# In[7]:


def extract_text_from_log(file_path):
    text=[]
    try:
        with open(file_path, "r", encoding="utf-8") as file:
            for line in file:
                text.append(line)
    except:
        text=""
        error_files.append(file_path)
    return "".join(text)


# In[8]:


def index_pdfs():
    folder_selected = filedialog.askdirectory()
    DirList=[]
    query = entry_words.get().lower()
    if folder_selected:
        entry_pdf.delete(0, tk.END)
        entry_pdf.insert(0, folder_selected)
    folder_path = entry_pdf.get()
    for dirname in os.scandir(folder_selected):
        if dirname.is_dir():
            DirList.append({"DirName":dirname.name,"DirPath":dirname.path})
    for j in DirList:
        j["DirPath"] = j["DirPath"].replace("/","\\")
    total_files = 0
    for i in DirList:
        total_files += len(os.listdir(i["DirPath"]))
    
    progress_bar["maximum"] = total_files
    progress_bar["value"] = 0
    processed_files = 0
    
    for i in DirList:
        for filename in os.listdir(i["DirPath"]):
            file_path = os.path.join(i["DirPath"], filename)
            current_file_label.config(text=f"{file_path}")
            root.update_idletasks()
            
            if filename.endswith(".pdf"):
                text = extract_text_from_pdf(file_path)
                text = text.lower()
                if query in text:
                    found_files.append(filename)
                    ThePath.append(file_path.replace("\\","/"))
                ##file_data.append({"path": file_path, "filename": filename, "content": text})
                del text
            elif filename.endswith(".pptx"):
                text = extract_text_from_pptx(file_path)
                text = text.lower()
                if query in text:
                    found_files.append(filename)
                    ThePath.append(file_path.replace("\\","/"))
                ##file_data.append({"path": file_path, "filename": filename, "content": text})
                del text
            elif filename.endswith(".docx"):
                text = extract_text_from_docx(file_path)
                text = text.lower()
                if query in text:
                    found_files.append(filename)
                    ThePath.append(file_path.replace("\\","/"))
                ##file_data.append({"path": file_path, "filename": filename, "content": text})
                del text
            elif filename.endswith(".xlsx"):
                Flag = extract_text_from_xlsx(file_path)
                ##text = text.lower()
                if Flag:
                    found_files.append(filename)
                    ThePath.append(file_path.replace("\\","/"))
                ##file_data.append({"path": file_path, "filename": filename, "content": text})
            elif filename.endswith(".doc"):
                text = extract_text_from_docx(file_path)
                text = text.lower()
                if query in text:
                    found_files.append(filename)
                    ThePath.append(file_path.replace("\\","/"))
                ##file_data.append({"path": file_path, "filename": filename, "content": text})
            elif filename.endswith(".log"):
                ##file_path = os.path.join(i["DirPath"], filename)
                text = extract_text_from_log(file_path)
                text = text.lower()
                if query in text:
                    found_files.append(filename)
                    ThePath.append(file_path.replace("\\","/"))
                ##file_data.append({"path": file_path, "filename": filename, "content":text})
                del text
            
            processed_files += 1
            progress_bar["value"] = processed_files
            root.update_idletasks()
    if found_files:
        ##messagebox.showinfo("Success", "Files indexed successfully!")
        messagebox.showinfo("Success", "Files found successfully!")
    else:
        messagebox.showwarning("Warning", "Error Occured!")


# In[9]:


def save_as():
    ##query = entry_words.get()
    temp=""
    if not found_files:
        messagebox.showinfo("Results", "The word can't be found!")
        return

    #messagebox.showinfo("Results", "Word found in:\n" + "\n".join(found_files))
    save_path = entry_output.get()
    query = entry_words.get()
    pdfmetrics.registerFont(TTFont('NotoFull', "C:/Users/B103040059/Downloads/PDFSearcher/NotoSansTC-VariableFont_wght.ttf"))
    if save_path:
        c = canvas.Canvas(save_path, pagesize=A4)
        width, height = A4
        Height = 20
        text_y = height - inch
        c.setFont("NotoFull", 12)
        c.drawString(100, text_y + 20, f"Click the links below to open the files with the word {query}:")
        text_y -= Height
        for i in range(len(found_files)):
            c.drawString(100, text_y, found_files[i])
            c.linkURL("file:///"+quote(ThePath[i]), (100, text_y - 2, 300, text_y + 12), relative=0)
            text_y -= Height
            if text_y < inch:
                c.showPage()
                c.setFont("NotoFull", 12)
                text_y = height - inch
        # error section
        text_y -= Height
        if text_y < inch:
            c.showPage()
            c.setFont("NotoFull", 12)
            text_y = height - inch
        c.drawString(100, text_y + 20, "Below the files that can't be opened or corrupted:")
        text_y -= Height
        for i in range(len(error_files)):
            c.drawString(100, text_y, error_files[i][56:])
            text_y -= Height
            if text_y < inch:
                c.showPage()
                c.setFont("NotoFull", 12)
                text_y = height - inch
        c.save()
        messagebox.showinfo("Saved", "Results saved successfully!")
    else:
        messagebox.showwarning("Warning", "Error saving file!")


# In[10]:


def select_output():
    file_path = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF Files", "*.pdf")])
    if file_path:
        entry_output.delete(0, tk.END)
        entry_output.insert(0, file_path)


# In[11]:


root = tk.Tk()
root.title("PDF Word Searcher")
root.geometry("600x400")

root.grid_columnconfigure(0, weight=1)
file_data = []
found_files = []
ThePath = []
error_files=[]
Flag_doc=0

tk.Label(root, text="Enter words to search:").pack()
entry_words = tk.Entry(root, width=50)
entry_words.pack()

tk.Label(root, text="Select Folder:").pack()
entry_pdf = tk.Entry(root, width=50) ##variable to pass the user input to the function
entry_pdf.pack()
tk.Button(root, text="Browse", command=index_pdfs).pack()

progress_bar = ttk.Progressbar(root, orient="horizontal", length=300, mode="determinate")
progress_bar.pack(pady=5)

current_file_label = tk.Label(root, text="", anchor="w", justify="left")
current_file_label.pack(pady=10)

tk.Label(root, text="Save As:").pack()
entry_output = tk.Entry(root, width=50)
entry_output.pack()
tk.Button(root, text="Browse", command=select_output).pack()

tk.Button(root, text="Generate List file", command=save_as).pack()
root.mainloop()

